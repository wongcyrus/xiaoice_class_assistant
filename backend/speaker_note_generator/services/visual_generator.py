"""Visual generation service for speaker note generator."""

import io
import logging
import os
from typing import Optional

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from google.adk.agents import LlmAgent

from utils.agent_utils import run_visual_agent

logger = logging.getLogger(__name__)


class VisualGenerator:
    """Service for generating enhanced slide visuals."""

    def __init__(
        self,
        designer_agent: LlmAgent,
        output_dir: str,
        skip_generation: bool = False,
    ):
        """
        Initialize the visual generator.

        Args:
            designer_agent: Agent for generating slide designs
            output_dir: Directory to save generated visuals
            skip_generation: Whether to skip visual generation
        """
        self.designer_agent = designer_agent
        self.output_dir = output_dir
        self.skip_generation = skip_generation
        self.previous_image: Optional[Image.Image] = None

        # Ensure output directory exists
        if not skip_generation:
            os.makedirs(output_dir, exist_ok=True)

    async def generate_visual(
        self,
        slide_idx: int,
        slide_image: Image.Image,
        speaker_notes: str,
        retry_errors: bool = False,
    ) -> Optional[bytes]:
        """
        Generate an enhanced visual for a slide.

        Args:
            slide_idx: Slide index (1-based)
            slide_image: Original slide image
            speaker_notes: Generated speaker notes
            retry_errors: Whether to regenerate existing images

        Returns:
            Image bytes if generated, None otherwise
        """
        if self.skip_generation:
            logger.info(
                f"Skipping visual generation for Slide {slide_idx} "
                "(skip-visuals active)."
            )
            return None

        # Check if visual already exists
        img_filename = f"slide_{slide_idx}_reimagined.png"
        img_path = os.path.join(self.output_dir, img_filename)

        if os.path.exists(img_path) and not retry_errors:
            logger.info(
                f"Visual already exists for Slide {slide_idx} "
                f"({img_filename}). Skipping generation."
            )
            try:
                with open(img_path, "rb") as f:
                    img_bytes = f.read()
                self._update_previous_image(img_bytes)
                return img_bytes
            except Exception as e:
                logger.error(f"Failed to load existing image {img_path}: {e}")
                # Fall through to regenerate

        # Generate new visual
        logger.info(f"--- Generating Visual for Slide {slide_idx} ---")

        logo_instruction = self._get_logo_instruction(slide_idx)
        designer_prompt = self._build_designer_prompt(
            speaker_notes,
            logo_instruction
        )

        designer_images = [slide_image]
        if self.previous_image:
            designer_images.append(self.previous_image)

        img_bytes = await run_visual_agent(
            self.designer_agent,
            designer_prompt,
            images=designer_images
        )

        if img_bytes:
            try:
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                logger.info(f"Saved reimagined slide to: {img_path}")
                self._update_previous_image(img_bytes)
            except Exception as e:
                logger.error(f"Failed to save generated image: {e}")
        else:
            logger.warning(f"No image generated for Slide {slide_idx}")
            self.previous_image = None

        return img_bytes

    def replace_slide_with_visual(
        self,
        prs,
        slide,
        img_path: str,
        speaker_notes: str,
    ) -> bool:
        """
        Replace slide content with generated visual and add notes to notes section.

        Args:
            prs: PowerPoint presentation object
            slide: PowerPoint slide object to modify
            img_path: Path to the generated image
            speaker_notes: Speaker notes to add to the notes section

        Returns:
            True if successful, False otherwise
        """
        try:
            # Remove all shapes from the slide
            for shape in list(slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)

            # Get slide dimensions from the presentation
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # Add the reimagined image to fill the entire slide
            slide.shapes.add_picture(
                img_path,
                left=0,
                top=0,
                width=slide_width,
                height=slide_height
            )

            # Add speaker notes to the notes section
            if not slide.has_notes_slide:
                slide.notes_slide
            slide.notes_slide.notes_text_frame.text = speaker_notes

            logger.info(
                f"Replaced slide content with reimagined visual."
            )
            return True

        except Exception as e:
            logger.error(f"Failed to replace slide with visual: {e}")
            return False

    def add_visual_to_presentation(
        self,
        prs: Presentation,
        slide_idx: int,
        img_path: str,
        speaker_notes: str,
    ) -> bool:
        """
        Add a generated visual as a new slide in the presentation.

        Args:
            prs: PowerPoint presentation object
            slide_idx: Original slide index
            img_path: Path to the generated image
            speaker_notes: Speaker notes to add to the slide

        Returns:
            True if successful, False otherwise
        """
        try:
            # Find a blank layout
            try:
                blank_layout = prs.slide_layouts[6]  # Usually blank
            except IndexError:
                logger.warning(
                    "Could not find blank slide layout (index 6), "
                    "using first available."
                )
                blank_layout = prs.slide_layouts[0]

            # Add new slide
            new_slide = prs.slides.add_slide(blank_layout)

            # Embed the generated image
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(5)
            new_slide.shapes.add_picture(
                img_path, left, top,
                width=width, height=height
            )

            # Add speaker notes as text
            txBox = new_slide.shapes.add_textbox(
                Inches(0.5), Inches(5.7),
                Inches(9), Inches(1.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Generated Notes for Slide {slide_idx}:\n{speaker_notes}"
            p.font.size = Pt(10)

            logger.info(
                f"Added new slide with reimagined image and notes "
                f"for Slide {slide_idx}."
            )
            return True

        except Exception as e:
            logger.error(f"Failed to add reimagined slide to PPTX: {e}")
            return False

    def _get_logo_instruction(self, slide_idx: int) -> str:
        """Get logo instruction based on slide position."""
        if slide_idx == 1:
            return (
                "You MUST prominently feature the logo/branding from "
                "IMAGE 1 (Original Draft Slide) in an appropriate corner."
            )
        else:
            return (
                "DO NOT include any logos or branding elements. "
                "Focus solely on content."
            )

    def _build_designer_prompt(
        self,
        speaker_notes: str,
        logo_instruction: str
    ) -> str:
        """Build the prompt for the designer agent."""
        style_ref = (
            "Style Reference (Previous Slide) provided."
            if self.previous_image
            else "N/A"
        )

        return (
            f"IMAGE 1: Original Slide Image provided.\n"
            f"IMAGE 2: {style_ref}\n"
            f"Speaker Notes: \"{speaker_notes}\"\n\n"
            f"TASK: Generate the high-fidelity slide image now.\n"
            f"CONTEXT: {logo_instruction}\n"
        )

    def _update_previous_image(self, img_bytes: bytes) -> None:
        """Update the previous image for style consistency."""
        try:
            self.previous_image = Image.open(io.BytesIO(img_bytes))
        except Exception as e:
            logger.warning(
                f"Could not load generated image for next iteration: {e}"
            )
            self.previous_image = None

    def reset_style_context(self) -> None:
        """Reset the style context (previous image)."""
        self.previous_image = None
