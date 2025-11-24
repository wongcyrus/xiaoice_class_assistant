"""Presentation processing service for speaker note generator."""

import logging
import sys
import os
from typing import Optional, Dict, Any

import pymupdf
from PIL import Image
from pptx import Presentation
from google.genai import types
from google.adk.runners import InMemoryRunner
from google.adk.agents import LlmAgent
from google.adk.tools.agent_tool import AgentTool

from config import Config
from utils.agent_utils import run_stateless_agent
from utils.image_utils import register_image, unregister_image
from utils.progress_utils import (
    load_progress,
    save_progress,
    create_slide_key,
    get_progress_file_path,
    should_retry_errors,
)
from tools.agent_tools import AgentToolFactory
from services.visual_generator import VisualGenerator

logger = logging.getLogger(__name__)


class PresentationProcessor:
    """Main service for processing presentations and generating speaker notes."""

    def __init__(
        self,
        config: Config,
        supervisor_agent: LlmAgent,
        analyst_agent: LlmAgent,
        writer_agent: LlmAgent,
        auditor_agent: LlmAgent,
        overviewer_agent: LlmAgent,
        designer_agent: LlmAgent,
    ):
        """
        Initialize the presentation processor.

        Args:
            config: Configuration object
            supervisor_agent: Supervisor agent for orchestration
            analyst_agent: Agent for slide analysis
            writer_agent: Agent for writing speaker notes
            auditor_agent: Agent for auditing existing notes
            overviewer_agent: Agent for generating global context
            designer_agent: Agent for generating visuals
        """
        self.config = config
        self.supervisor_agent = supervisor_agent
        self.analyst_agent = analyst_agent
        self.writer_agent = writer_agent
        self.auditor_agent = auditor_agent
        self.overviewer_agent = overviewer_agent
        self.designer_agent = designer_agent

        # Initialize tool factory
        self.tool_factory = AgentToolFactory(
            analyst_agent=analyst_agent,
            writer_agent=writer_agent,
            auditor_agent=auditor_agent,
        )

        # Initialize visual generator
        self.visual_generator = VisualGenerator(
            designer_agent=designer_agent,
            output_dir=config.visuals_dir,
            skip_generation=config.skip_visuals,
        )

        # Progress tracking
        self.progress_file = get_progress_file_path(config.pptx_path)
        self.retry_errors = should_retry_errors()

        logger.info(f"Initialized processor with config: {config}")
        logger.info(f"Progress file: {self.progress_file}")
        logger.info(f"Retry errors: {self.retry_errors}")

    async def process(self) -> tuple[str, str]:
        """
        Process the presentation and generate speaker notes.

        Returns:
            Tuple of (notes_only_path, with_visuals_path)
        """
        logger.info(f"Processing PPTX: {self.config.pptx_path}")
        logger.info(f"Region: {os.environ.get('GOOGLE_CLOUD_LOCATION')}")

        # Load files - create two separate presentations
        prs_notes = Presentation(self.config.pptx_path)
        prs_visuals = Presentation(self.config.pptx_path)
        pdf_doc = pymupdf.open(self.config.pdf_path)
        limit = min(len(prs_notes.slides), len(pdf_doc))

        # Load progress
        progress = load_progress(self.progress_file)

        # Generate or load global context
        global_context = await self._get_global_context(
            pdf_doc, limit, progress
        )

        # Get presentation theme
        presentation_theme = self.config.get_presentation_theme()

        # Setup supervisor tools
        self._configure_supervisor_tools(presentation_theme, global_context)

        # Initialize supervisor runner
        supervisor_runner = await self._initialize_supervisor()

        # Process slides
        await self._process_slides(
            prs_notes, prs_visuals, pdf_doc, limit, progress,
            supervisor_runner, presentation_theme,
            global_context
        )

        # Save both presentations
        output_path_notes = self.config.output_path
        output_path_visuals = self.config.output_path_with_visuals

        prs_notes.save(output_path_notes)
        logger.info(f"Saved presentation with notes to: {output_path_notes}")

        prs_visuals.save(output_path_visuals)
        logger.info(f"Saved presentation with visuals to: {output_path_visuals}")

        return output_path_notes, output_path_visuals

    async def _get_global_context(
        self,
        pdf_doc,
        limit: int,
        progress: Dict[str, Any]
    ) -> str:
        """Generate or retrieve cached global context."""
        # Check if cached
        if (
            "global_context" in progress
            and progress["global_context"]
            and len(progress["global_context"]) > 50
            and not self.retry_errors
        ):
            logger.info("Using cached Global Context from progress file.")
            return progress["global_context"]

        # Generate new context
        logger.info("--- Pass 1: Generating Global Context ---")

        all_images = []
        for i in range(limit):
            pix = pdf_doc[i].get_pixmap(dpi=75)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            all_images.append(img)

        global_context = await run_stateless_agent(
            self.overviewer_agent,
            "Here are the slides for the entire presentation. Analyze them.",
            images=all_images
        )

        logger.info(f"Global Context Generated: {len(global_context)} chars")

        # Cache it
        progress["global_context"] = global_context
        save_progress(self.progress_file, progress)

        return global_context

    def _configure_supervisor_tools(
        self,
        presentation_theme: str,
        global_context: str
    ) -> None:
        """Configure the supervisor agent's tools."""
        self.supervisor_agent.tools = [
            AgentTool(agent=self.auditor_agent),
            self.tool_factory.create_analyst_tool(),
            self.tool_factory.create_writer_tool(
                presentation_theme,
                global_context
            ),
        ]

    async def _initialize_supervisor(self) -> InMemoryRunner:
        """Initialize and create supervisor session."""
        supervisor_runner = InMemoryRunner(
            agent=self.supervisor_agent,
            app_name="supervisor"
        )

        user_id = "supervisor_user"
        session_id = "supervisor_session"

        await supervisor_runner.session_service.create_session(
            app_name="supervisor",
            user_id=user_id,
            session_id=session_id
        )

        return supervisor_runner

    async def _process_slides(
        self,
        prs_notes: Presentation,
        prs_visuals: Presentation,
        pdf_doc,
        limit: int,
        progress: Dict[str, Any],
        supervisor_runner: InMemoryRunner,
        presentation_theme: str,
        global_context: str,
    ) -> None:
        """Process all slides in both presentations."""
        previous_slide_summary = "Start of presentation."
        user_id = "supervisor_user"
        session_id = "supervisor_session"

        for i in range(limit):
            slide_idx = i + 1
            slide_notes = prs_notes.slides[i]
            slide_visuals = prs_visuals.slides[i]
            pdf_page = pdf_doc[i]

            logger.info(f"--- Processing Slide {slide_idx} ---")

            # Get existing notes (from notes presentation)
            existing_notes = self._get_existing_notes(slide_notes)
            skey = create_slide_key(slide_idx, existing_notes)
            entry = progress["slides"].get(skey)

            # Register slide image
            image_id = f"slide_{slide_idx}"
            slide_image = self._extract_slide_image(pdf_page)
            register_image(image_id, slide_image)

            # Generate or retrieve speaker notes
            final_response, status = await self._process_slide_notes(
                slide_idx, image_id, existing_notes,
                previous_slide_summary, presentation_theme,
                global_context, entry, supervisor_runner,
                user_id, session_id
            )

            # Update slide notes in both presentations
            if status == "success":
                # Update notes in notes-only presentation
                self._update_slide_notes(slide_notes, final_response)
                
                # Update notes in visuals presentation
                self._update_slide_notes(slide_visuals, final_response)
                
                previous_slide_summary = final_response[:200]

                # Generate visual and replace slide in visuals presentation
                img_bytes = await self.visual_generator.generate_visual(
                    slide_idx, slide_image, final_response, self.retry_errors
                )

                if img_bytes:
                    img_path = os.path.join(
                        self.config.visuals_dir,
                        f"slide_{slide_idx}_reimagined.png"
                    )
                    # Replace the slide content with the visual
                    self.visual_generator.replace_slide_with_visual(
                        prs_visuals, slide_visuals, img_path, final_response
                    )

            # Update progress
            progress["slides"][skey] = {
                "slide_index": slide_idx,
                "existing_notes_hash": skey.split("_")[-1],
                "original_notes": existing_notes,
                "note": final_response,
                "status": status,
            }
            save_progress(self.progress_file, progress)

            # Cleanup
            unregister_image(image_id)

    def _get_existing_notes(self, slide) -> str:
        """Extract existing notes from a slide."""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            return slide.notes_slide.notes_text_frame.text.strip()
        return ""

    def _extract_slide_image(self, pdf_page) -> Image.Image:
        """Extract image from PDF page."""
        pix = pdf_page.get_pixmap(dpi=150)
        return Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

    async def _process_slide_notes(
        self,
        slide_idx: int,
        image_id: str,
        existing_notes: str,
        previous_slide_summary: str,
        presentation_theme: str,
        global_context: str,
        entry: Optional[Dict[str, Any]],
        supervisor_runner: InMemoryRunner,
        user_id: str,
        session_id: str,
    ) -> tuple[str, str]:
        """
        Process notes for a single slide.

        Returns:
            Tuple of (final_response, status)
        """
        # Check if already done
        if entry and entry.get("status") == "success" and not self.retry_errors:
            logger.info(f"Skipping generation for slide {slide_idx}")
            return entry.get("note", ""), "success"

        # Build supervisor prompt
        supervisor_prompt = self._build_supervisor_prompt(
            slide_idx, image_id, existing_notes,
            previous_slide_summary, presentation_theme,
            global_context
        )

        content = types.Content(
            role='user',
            parts=[types.Part.from_text(text=supervisor_prompt)]
        )

        # Run supervisor
        final_response = ""
        status = "pending"

        try:
            for event in supervisor_runner.run(
                user_id=user_id,
                session_id=session_id,
                new_message=content,
            ):
                if getattr(event, "content", None) and event.content.parts:
                    for part in event.content.parts:
                        fn_call = getattr(part, "function_call", None)
                        if fn_call:
                            print(
                                f"\n[Supervisor] 📞 calling tool: {fn_call.name}"
                            )
                        text = getattr(part, "text", "") or ""
                        final_response += text
        except Exception as e:
            logger.error(f"Error in supervisor loop: {e}")
            status = "error"
            return final_response, status

        # Process result
        final_response = final_response.strip()

        if not final_response:
            # Try fallback to last writer output
            last_output = self.tool_factory.last_writer_output
            if last_output:
                logger.info(
                    f"Supervisor returned empty text, using fallback content "
                    f"({len(last_output)} chars)."
                )
                final_response = last_output
                status = "success"
                self.tool_factory.reset_writer_output()
            else:
                logger.warning(
                    f"Supervisor loop finished with empty response "
                    f"for Slide {slide_idx}."
                )
                status = "error"
        else:
            status = "success"
            self.tool_factory.reset_writer_output()

        return final_response, status

    def _build_supervisor_prompt(
        self,
        slide_idx: int,
        image_id: str,
        existing_notes: str,
        previous_slide_summary: str,
        presentation_theme: str,
        global_context: str,
    ) -> str:
        """Build the prompt for the supervisor agent."""
        return (
            f"Here is Slide {slide_idx}.\n"
            f"Existing Notes: \"{existing_notes}\"\n"
            f"Image ID: \"{image_id}\"\n"
            f"Previous Slide Summary: \"{previous_slide_summary}\"\n"
            f"Theme: \"{presentation_theme}\"\n"
            f"Global Context: \"{global_context}\"\n\n"
            f"Please proceed with the workflow."
        )

    def _update_slide_notes(self, slide, notes: str) -> None:
        """Update the notes on a slide."""
        try:
            if not slide.has_notes_slide:
                slide.notes_slide
            slide.notes_slide.notes_text_frame.text = notes
        except Exception as e:
            logger.error(f"Failed to update slide notes: {e}")
