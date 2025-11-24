#!/usr/bin/env python3
"""
Speaker Note Generator
Enhances PowerPoint presentations by generating speaker notes using a Supervisor-Tool Multi-Agent System.

Refactored to use Python-based Agent definitions and support Image Generation.
"""

import argparse
import asyncio
import logging
import os
import sys
import io
import json
import hashlib
import tempfile
from typing import Dict, Any, List, Optional

import pymupdf  # fitz
from PIL import Image
from pptx import Presentation

from google.adk.runners import InMemoryRunner
from google.genai import types
from google.adk.agents import LlmAgent

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(levelname)s:%(name)s:%(message)s',
    stream=sys.stdout
)
logger = logging.getLogger(__name__)

# Global registry for images
IMAGE_REGISTRY: Dict[str, Image.Image] = {}

def _create_image_part(image: Image.Image) -> types.Part:
    """Helper to create a Part object from a PIL Image safely."""
    if hasattr(types.Part, 'from_image'):
        return types.Part.from_image(image=image)
    else:
        buf = io.BytesIO()
        image.save(buf, format='PNG')
        img_bytes = buf.getvalue()
        if hasattr(types.Part, 'from_bytes'):
            return types.Part.from_bytes(
                data=img_bytes,
                mime_type='image/png'
            )
        else:
            return types.Part(
                mime_type='image/png',
                data=img_bytes
            )

async def run_stateless_agent(
    agent: LlmAgent,
    prompt: str,
    images: List[Image.Image] = None,
) -> str:
    """Helper to run a stateless single-turn agent (Text Output Only)."""
    runner = InMemoryRunner(agent=agent, app_name=agent.name)
    user_id = "system_user"
    
    parts = [types.Part.from_text(text=prompt)]
    if images:
        for img in images:
            try:
                parts.append(_create_image_part(img))
            except Exception as e:
                logger.error(f"Failed to attach image part: {e}")

    content = types.Content(role='user', parts=parts)
    
    # Create new session for statelessness
    session = await runner.session_service.create_session(
        app_name=agent.name,
        user_id=user_id,
    )
    
    # Safely determine session id
    resolved_session_id = f"session_{user_id}" # Default fallback
    try:
        if hasattr(session, "session_id"): resolved_session_id = session.session_id
        elif hasattr(session, "id"): resolved_session_id = session.id
    except:
        pass

    print(f"\n┌── [Agent: {agent.name}]")
    print(f"│ Task: {prompt.strip()[:500].replace(chr(10), ' ') + ('...' if len(prompt) > 500 else '')}")
    if images:
        print(f"│ [{len(images)} Images Attached]")

    response_text = ""
    try:
        # Run agent
        for event in runner.run(
            user_id=user_id,
            session_id=resolved_session_id,
            new_message=content,
        ):
            if getattr(event, "content", None) and event.content.parts:
                for part in event.content.parts:
                     # Only extract text parts
                     txt = getattr(part, "text", "") or ""
                     response_text += txt
            else:
                 # Debug: Print raw event if no content parts found
                 # print(f"DEBUG: Event with no content parts: {event}")
                 pass
    except Exception as e:
        logger.error(f"Error running agent {agent.name}: {e}")
        return f"Error: {e}"
    
    if not response_text.strip():
        print(f"⚠️ WARNING: Agent {agent.name} returned empty text.")
    
    print(f"└-> Response: {response_text.strip()[:500].replace(chr(10), ' ') + ('...' if len(response_text) > 500 else '')}")
    return response_text.strip()

async def run_visual_agent(
    agent: LlmAgent,
    prompt: str,
    images: List[Image.Image] = None,
) -> Optional[bytes]:
    """Helper to run a stateless agent and capture generated IMAGE."""
    runner = InMemoryRunner(agent=agent, app_name=agent.name)
    user_id = "system_user"
    
    parts = [types.Part.from_text(text=prompt)]
    if images:
        for img in images:
            try:
                parts.append(_create_image_part(img))
            except Exception as e:
                logger.error(f"Failed to attach image part: {e}")

    content = types.Content(role='user', parts=parts)
    
    session = await runner.session_service.create_session(
        app_name=agent.name,
        user_id=user_id,
    )
    # Safely determine session id
    resolved_session_id = f"session_{user_id}" 
    try:
        if hasattr(session, "session_id"): resolved_session_id = session.session_id
        elif hasattr(session, "id"): resolved_session_id = session.id
    except:
        pass

    print(f"\n┌── [Agent: {agent.name} (Visual)]")
    print(f"│ Task: {prompt.strip()[:500].replace(chr(10), ' ') + ('...' if len(prompt) > 500 else '')}")

    generated_image_bytes = None

    try:
        for event in runner.run(
            user_id=user_id,
            session_id=resolved_session_id,
            new_message=content,
        ):
            if getattr(event, "content", None) and event.content.parts:
                for part in event.content.parts:
                    # Check for inline_data or file_data
                    inline_data = getattr(part, "inline_data", None)
                    if inline_data:
                        generated_image_bytes = inline_data.data
                        print("│ [Received Image Data]")
                    # Some libraries might use specific image attributes
                    elif hasattr(part, "image"):
                        # Handle if it's already an object
                        pass
    except Exception as e:
        logger.error(f"Error running visual agent {agent.name}: {e}")
        return None

    if generated_image_bytes:
        print(f"└-> Response: [Image Generated ({len(generated_image_bytes)} bytes)]")
    else:
        print(f"└-> Response: [No Image Generated]")
        
    return generated_image_bytes

async def process_presentation(
    pptx_path: str,
    pdf_path: str,
    course_id: str = None,
):
    # Late Import to allow env var configuration
    sys.path.append(os.path.dirname(os.path.abspath(__file__)))
    from google.adk.tools.agent_tool import AgentTool # Import AgentTool
    from agents.supervisor import supervisor_agent
    from agents.analyst import analyst_agent
    from agents.overviewer import overviewer_agent
    from agents.designer import designer_agent
    from agents.writer import writer_agent
    from agents.auditor import auditor_agent # Import auditor_agent

    # Tool wrapper for Analyst
    async def call_analyst(image_id: str) -> str:
        """Tool: Analyzes the slide image."""
        logger.info(f"[Tool] call_analyst invoked for image_id: {image_id}")
        image = IMAGE_REGISTRY.get(image_id)
        if not image:
            return "Error: Image not found."
        
        prompt_text = "Analyze this slide image."
        return await run_stateless_agent(analyst_agent, prompt_text, images=[image])

    async def speech_writer(
        analysis: str,
        previous_context: str,
        theme: str,
        global_context: str = "No global context provided." 
    ) -> str:
        """Tool: Writes the script."""
        logger.info("[Tool] speech_writer invoked.")
        prompt = (
            f"SLIDE_ANALYSIS:\n{analysis}\n\n"
            f"PRESENTATION_THEME: {theme}\n"
            f"PREVIOUS_CONTEXT: {previous_context}\n"
            f"GLOBAL_CONTEXT: {global_context}\n"
        )
        result = await run_stateless_agent(writer_agent, prompt)
        if not result or not result.strip():
            logger.warning("[Tool] speech_writer returned empty text. Returning fallback.")
            return "Error: The writer agent failed to generate a script. Please try again or use a placeholder."
        return result
    
    logger.info(f"Processing PPTX: {pptx_path}")
    logger.info(f"Region: {os.environ.get('GOOGLE_CLOUD_LOCATION')}")
    
    # Load files
    prs = Presentation(pptx_path)
    pdf_doc = pymupdf.open(pdf_path)
    limit = min(len(prs.slides), len(pdf_doc))

    progress = load_progress(progress_file)

    # Check if global_context already exists in progress file
    if "global_context" in progress and progress["global_context"] and len(progress["global_context"]) > 50 and not retry_errors:
        global_context = progress["global_context"]
        logger.info("Using cached Global Context from progress file.")
    else:
        logger.info("--- Pass 1: Generating Global Context ---")
        all_images = []
        for i in range(limit):
            pix = pdf_doc[i].get_pixmap(dpi=75) 
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            all_images.append(img)
        
        global_context = await run_stateless_agent(
            overviewer_agent,
            "Here are the slides for the entire presentation. Analyze them.",
            images=all_images
        )
        logger.info(f"Global Context Generated: {len(global_context)} chars")
        
        progress["global_context"] = global_context
        save_progress(progress_file, progress)

    # 1. Pass 2: Slide Loop
    previous_reimagined_image: Optional[Image.Image] = None

    for i in range(limit):
        slide_idx = i + 1
        slide = prs.slides[i]
        pdf_page = pdf_doc[i]
        
        logger.info(f"--- Processing Slide {slide_idx} ---")
        
        existing_notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            existing_notes = slide.notes_slide.notes_text_frame.text.strip()

        skey = slide_key(slide_idx, existing_notes)
        entry = progress["slides"].get(skey)

        # Register Image
        image_id = f"slide_{slide_idx}"
        pix = pdf_page.get_pixmap(dpi=150)
        slide_image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        IMAGE_REGISTRY[image_id] = slide_image

        final_response = ""
        status = "pending"

        # Skip generation if done, BUT check if visual is done too?
        # For now, assume if notes are done, we reuse them.
        if entry and entry.get("status") == "success" and not retry_errors:
            final_response = entry.get("note", "")
            logger.info(f"Skipping generation for slide {slide_idx}")
            status = "success"
        else:
            # 2. Prompt Supervisor
            supervisor_prompt = (
                f"Here is Slide {slide_idx}.\n"
                f"Existing Notes: \"{existing_notes}\"\n"
                f"Image ID: \"{image_id}\"\n"
                f"Previous Slide Summary: \"{previous_slide_summary}\"\n"
                f"Theme: \"{presentation_theme}\"\n"
                f"Global Context: \"{global_context}\"\n\n"
                f"Please proceed with the workflow."
            )

            content = types.Content(
                role='user',
                parts=[types.Part.from_text(text=supervisor_prompt)]
            )

            try:
                for event in supervisor_runner.run(
                    user_id=user_id,
                    session_id=session_id,
                    new_message=content,
                ):
                    if getattr(event, "content", None) and event.content.parts:
                        for part in event.content.parts:
                            fn_call = getattr(part, "function_call", None)
                            if fn_call:
                                print(f"\n[Supervisor] 📞 calling tool: {fn_call.name}")
                            text = getattr(part, "text", "") or ""
                            final_response += text
            except Exception as e:
                logger.error(f"Error in supervisor loop: {e}")
                status = "error"

            final_response = final_response.strip()
            if not final_response: 
                # Fallback: If writer was called, maybe we can extract the thought process or retry
                # For now, let's mark as error but log heavily
                logger.warning(f"Supervisor loop finished with empty response for Slide {slide_idx}. Tool calls might have succeeded but text was lost.")
                status = "error"
            else: status = "success"

        # 3. Update Notes
        if status == "success":
            try:
                if not slide.has_notes_slide: slide.notes_slide
                slide.notes_slide.notes_text_frame.text = final_response
            except: pass
            
            previous_slide_summary = final_response[:200]

            # 4. VISUALIZATION STEP (New)
            logger.info(f"--- Processing Visual for Slide {slide_idx} ---")
            
            # Create a hash of the generated notes for unique image filename
            notes_hash = hashlib.sha256(final_response.encode("utf-8")).hexdigest()[:8]
            vis_dir = os.path.join(os.path.dirname(pptx_path), "visuals")
            os.makedirs(vis_dir, exist_ok=True)
            img_filename = f"slide_{slide_idx}_{notes_hash}_reimagined.png"
            img_path = os.path.join(vis_dir, img_filename)

            img_bytes = None

            if os.path.exists(img_path) and not retry_errors: # Skip if exists and not retrying errors
                logger.info(f"Visual already exists for Slide {slide_idx} ({img_filename}). Skipping generation.")
                try:
                    with open(img_path, "rb") as f:
                        img_bytes = f.read()
                    # Update Previous Image for next iteration (Style Consistency)
                    try:
                         previous_reimagined_image = Image.open(io.BytesIO(img_bytes))
                    except Exception as img_err:
                         logger.warning(f"Could not load generated image for next iteration context: {img_err}")
                         previous_reimagined_image = None
                except Exception as e:
                    logger.error(f"Failed to load existing image {img_path}: {e}")
                    img_bytes = None # Force re-generation if load fails
            
            if img_bytes is None: # Only generate if not skipped or load failed
                logger.info(f"--- Generating Visual for Slide {slide_idx} ---")
                
                logo_instruction = ""
                if slide_idx == 1:
                    logo_instruction = "You MUST prominently feature the logo/branding from IMAGE 1 (Original Draft Slide) in an appropriate corner."
                else:
                    logo_instruction = "DO NOT include any logos or branding elements. Focus solely on content."

                designer_prompt_text = (
                    f"IMAGE 1: Original Slide Image provided.\n"
                    f"IMAGE 2: {'Style Reference (Previous Slide) provided.' if previous_reimagined_image else 'N/A'}\n"
                    f"Speaker Notes: \"{final_response}\"\n\n"
                    f"TASK: Generate the high-fidelity slide image now.\n"
                    f"CONTEXT: {logo_instruction}\n"
                )
                
                designer_images = [slide_image]
                if previous_reimagined_image:
                    designer_images.append(previous_reimagined_image)

                img_bytes = await run_visual_agent(
                    designer_agent,
                    designer_prompt_text,
                    images=designer_images
                )
            
            if img_bytes:
                # Save to disk if newly generated or re-saved for consistency
                try:
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                    logger.info(f"Saved reimagined slide to: {img_path}")
                    
                    # Update Previous Image for next iteration (Style Consistency)
                    try:
                         previous_reimagined_image = Image.open(io.BytesIO(img_bytes))
                    except Exception as img_err:
                         logger.warning(f"Could not load generated image for next iteration context: {img_err}")
                         previous_reimagined_image = None

                    # Create a new slide in the presentation to embed the image and notes
                    # Find a blank layout (usually slide_layouts[6])
                    try:
                        blank_slide_layout = prs.slide_layouts[6] # Usually the blank layout
                    except IndexError:
                        logger.warning("Could not find blank slide layout (index 6), using first available.")
                        blank_slide_layout = prs.slide_layouts[0] # Fallback

                    # Add a new slide to the presentation
                    new_slide = prs.slides.add_slide(blank_slide_layout)
                    
                    from pptx.util import Inches, Pt # Import here to ensure it's available

                    # Embed the generated image
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(5)
                    new_slide.shapes.add_picture(img_path, left, top, width=width, height=height)

                    # Add the speaker notes as text on the new slide
                    txBox = new_slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(9), Inches(1.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = f"Generated Notes for Slide {slide_idx}:\n{final_response}"
                    p.font.size = Pt(10)

                    logger.info(f"Added new slide with reimagined image and notes for Slide {slide_idx}.")

                except Exception as e:
                    logger.error(f"Failed to add reimagined slide to PPTX: {e}")
            else:
                logger.warning(f"No image generated for Slide {slide_idx}")
                previous_reimagined_image = None


        # Update progress
        progress["slides"][skey] = {
            "slide_index": slide_idx,
            "existing_notes_hash": skey.split("_")[-1],
            "original_notes": existing_notes,
            "note": final_response,
            "status": status,
        }
        save_progress(progress_file, progress)
        del IMAGE_REGISTRY[image_id]

    # Save
    output_path = pptx_path.replace(".pptx", "_enhanced.pptx")
    prs.save(output_path)
    logger.info(f"Saved enhanced deck to: {output_path}")

def main():
    parser = argparse.ArgumentParser(description="Generate Speaker Notes with Supervisor Agent")
    parser.add_argument("--pptx", required=True, help="Path to input PPTX")
    parser.add_argument("--pdf", required=True, help="Path to input PDF")
    parser.add_argument("--course-id", help="Optional: Course ID to fetch theme context")
    parser.add_argument("--progress-file", help="Override path for progress JSON file")
    parser.add_argument("--retry-errors", action="store_true", help="Retry slides previously marked as error")
    parser.add_argument("--region", help="Google Cloud Region (default: global)", default="global")

    args = parser.parse_args()

    if not os.path.exists(args.pptx) or not os.path.exists(args.pdf):
        print("Error: Input files not found.")
        return

    if args.progress_file:
        os.environ["SPEAKER_NOTE_PROGRESS_FILE"] = args.progress_file
    if args.retry_errors:
        os.environ["SPEAKER_NOTE_RETRY_ERRORS"] = "true"
    
    # Set Google Cloud Location based on arg (override env if provided)
    if args.region:
        os.environ["GOOGLE_CLOUD_LOCATION"] = args.region
    elif "GOOGLE_CLOUD_LOCATION" not in os.environ:
        os.environ["GOOGLE_CLOUD_LOCATION"] = "global"

    asyncio.run(process_presentation(args.pptx, args.pdf, args.course_id))

if __name__ == "__main__":
    main()
