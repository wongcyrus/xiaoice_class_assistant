#!/usr/bin/env python3
"""
Pre-generate presentation_messages from PPTX speaker notes and cache them
for fast retrieval by generate_presentation_message().

This tool:
1. Reads speaker notes from each slide
2. Uses Gemini to generate a message for each slide
3. Caches each message by speaker notes content (not slide number)
4. Pre-generates speech files and uploads to GCS bucket

Usage:
  python preload_presentation_messages.py \
    --pptx /path/to/deck.pptx \
    --languages en,zh

Notes:
  - Generates ONE message per unique speaker notes content
  - Cache key format: "v1:{language}:{hash(speaker_notes)}"
  - No slide numbers in cache - works even if slides reordered
  - VBA sends current slide's speaker notes for cache lookup
  - Duplicate speaker notes (same content) share same cache entry
  - Speech files named: speech_{lang}_{content_hash}.mp3
  - Bucket name read from config.py (generated from Terraform outputs)
"""

import argparse
import asyncio
import hashlib
import logging
import os
import sys
from typing import Dict, List

from google.cloud import firestore, texttospeech, storage
from google.adk.agents import config_agent_utils
from google.adk.runners import InMemoryRunner
from google.genai import types
from pptx import Presentation

# Add path to backend/functions/config to import course_utils
sys.path.append(os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "functions", "config"))
import course_utils

# Import config for bucket name
try:
    import config
except ImportError:
    config = None

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(levelname)s:%(name)s:%(message)s',
    stream=sys.stdout
)
logger = logging.getLogger(__name__)


def create_agent():
    """Create and return an ADK agent from YAML config."""
    # Get path to the agent config file
    script_dir = os.path.dirname(os.path.abspath(__file__))
    backend_dir = os.path.dirname(script_dir)
    config_file_path = os.path.join(
        backend_dir,
        "functions",
        "config",
        "presenter_agent",
        "root_agent.yaml"
    )
    
    # Load the agent from the config file using utility function
    return config_agent_utils.from_config(config_file_path)


def parse_languages(s: str) -> List[str]:
    return [x.strip() for x in s.split(",") if x.strip()]


def _normalize_context(context: str) -> str:
    """Normalize context by trimming and collapsing whitespace."""
    if not context:
        return ""
    # Collapse all whitespace runs to a single space and strip ends
    return " ".join(str(context).split())


def _cache_key(language_code: str, context: str) -> str:
    """Build a stable, short cache key safe for Firestore doc IDs.

    Uses lowercase language + 12-char SHA256 of normalized context.
    Avoids very long document IDs and ensures consistent lookups.
    Must match the logic in firestore_utils.py
    """
    lang = (language_code or "").strip().lower() or "unknown"
    norm_ctx = _normalize_context(context)
    if not norm_ctx:
        return f"v1:{lang}:default"
    digest = hashlib.sha256(norm_ctx.encode("utf-8")).hexdigest()[:12]
    return f"v1:{lang}:{digest}"


def generate_message_from_notes(
    runner: InMemoryRunner,
    speaker_notes: str,
    language_code: str,
    slide_num: int
) -> str:
    """Generate a presentation message from speaker notes using Gemini.
    
    Args:
        runner: ADK runner instance
        speaker_notes: Raw speaker notes from all slides
        language_code: Target language (e.g., 'en', 'zh')
        slide_num: Slide number for logging (0 if combined notes)
        
    Returns:
        Generated message text
    """
    prompt = (
        f"Transform the following presentation speaker notes into a "
        f"clear, engaging message for students in {language_code} "
        f"language.\\n\\n"
        f"Speaker Notes:\\n{speaker_notes}\\n\\n"
        f"Generate a concise message (2-4 sentences) that captures the "
        f"key points. Return ONLY the message text, no explanations."
    )

    try:
        session_id = f"preload_slide{slide_num}_{language_code}"
        user_id = "preload_system"

        # Get or create session
        session = asyncio.run(
            runner.session_service.get_session(
                app_name='presentation_preloader',
                user_id=user_id,
                session_id=session_id,
            )
        )
        if session is None:
            session = asyncio.run(
                runner.session_service.create_session(
                    app_name='presentation_preloader',
                    user_id=user_id,
                    session_id=session_id,
                )
            )

        content = types.Content(
            role='user',
            parts=[types.Part.from_text(text=prompt)]
        )

        generated_text = ""
        for event in runner.run(
            user_id=user_id,
            session_id=session_id,
            new_message=content,
        ):
            if getattr(event, "content", None) and event.content.parts:
                part0 = event.content.parts[0]
                text = getattr(part0, "text", "") or ""
                if text:
                    generated_text += text

        result = generated_text.strip()
        return result
    except Exception as e:
        logger.exception(
            "Failed to generate message for slide %d: %s",
            slide_num,
            e
        )
        return ""


def get_speaker_notes(prs: Presentation, slide_index: int) -> str:
    """Extract speaker notes from a slide.
    
    Args:
        prs: PowerPoint presentation object
        slide_index: 1-based slide index
        
    Returns:
        Speaker notes text, or empty string if no notes
    """
    if slide_index < 1 or slide_index > len(prs.slides):
        raise IndexError(
            f"slide index {slide_index} out of range (1..{len(prs.slides)})"
        )
    slide = prs.slides[slide_index - 1]
    
    # Access speaker notes
    if not slide.has_notes_slide:
        return ""
    
    notes_slide = slide.notes_slide
    if not notes_slide or not notes_slide.notes_text_frame:
        return ""
    
    return (notes_slide.notes_text_frame.text or "").strip()


def cache_message(
    db: firestore.Client,
    language_code: str,
    context: str,
    message: str,
    course_id: str = None
):
    """Write a cache entry for generate_presentation_message lookup."""
    cache_key = _cache_key(language_code, context)
    norm_ctx = _normalize_context(context)
    cache_ref = db.collection("langbridge_presentation_cache").document(cache_key)
    
    data = {
        "message": message,
        "language_code": (language_code or "").strip().lower(),
        "context": norm_ctx,
        "context_hash": cache_key.rsplit(":", 1)[-1],
        "updated_at": firestore.SERVER_TIMESTAMP
    }
    
    if course_id:
        data["course_ids"] = firestore.ArrayUnion([course_id])
        
    cache_ref.set(data, merge=True)


def update_config(db: firestore.Client, messages: Dict[str, str]):
    """Update presentation_messages in config document."""
    doc_ref = db.collection("langbridge_config").document("messages")
    doc = doc_ref.get()
    current = doc.to_dict() if doc.exists else {}
    merged = dict(current or {})
    merged["presentation_messages"] = messages
    merged["updated_at"] = firestore.SERVER_TIMESTAMP
    doc_ref.set(merged)


def generate_speech_file(
    bucket_name: str,
    message: str,
    language_code: str,
    voice_params: texttospeech.VoiceSelectionParams = None
) -> str:
    """Generate speech file and upload to bucket.
    
    Returns:
        Filename of uploaded speech file
    """
    tts_client = texttospeech.TextToSpeechClient()
    storage_client = storage.Client()
    
    # Generate stable filename from message content and language
    content_hash = hashlib.sha256(
        f"{message}:{language_code}".encode("utf-8")
    ).hexdigest()[:12]
    filename = f"speech_{language_code}_{content_hash}.mp3"
    
    bucket = storage_client.bucket(bucket_name)
    blob = bucket.blob(filename)
    
    # Skip if already exists
    if blob.exists():
        logger.info("Speech file already exists: %s", filename)
        return filename
    
    # Determine voice params if not provided
    if not voice_params:
        if language_code.startswith("en"):
            voice_language = "en-US"
        elif language_code.startswith("zh"):
            voice_language = "zh-CN"
        else:
            voice_language = "en-US"
            
        voice_params = texttospeech.VoiceSelectionParams(
            language_code=voice_language,
            ssml_gender=texttospeech.SsmlVoiceGender.FEMALE
        )
    
    synthesis_input = texttospeech.SynthesisInput(text=message)
    audio_config = texttospeech.AudioConfig(
        audio_encoding=texttospeech.AudioEncoding.MP3,
        speaking_rate=1.0
    )
    
    tts_response = tts_client.synthesize_speech(
        input=synthesis_input,
        voice=voice_params,
        audio_config=audio_config
    )
    
    blob.upload_from_string(
        tts_response.audio_content,
        content_type="audio/mpeg"
    )
    logger.info("Generated speech file: %s", filename)
    return filename


def main():
    parser = argparse.ArgumentParser(
        description="Preload presentation message cache from PPTX "
                    "speaker notes"
    )
    parser.add_argument("--pptx", required=True, help="Path to PPTX file")
    parser.add_argument(
        "--languages", default="en", help="Comma list, e.g. en,zh"
    )
    parser.add_argument(
        "--course-id", help="Optional Course ID for config and cache tagging"
    )

    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        raise FileNotFoundError(args.pptx)
    
    # Get bucket name from config
    bucket_name = None
    if config and hasattr(config, 'speech_file_bucket'):
        bucket_name = config.speech_file_bucket
        if not bucket_name:
            logger.warning(
                "speech_file_bucket in config.py is empty. "
                "Run update_config_from_cdktf.sh first."
            )

    languages = parse_languages(args.languages)
    
    # If course_id is provided, override languages with course config
    if args.course_id:
        course_langs = course_utils.get_course_languages(args.course_id)
        if course_langs:
            languages = course_langs
            print(f"Using languages from course {args.course_id}: {', '.join(languages)}")
        else:
            print(f"Warning: Course {args.course_id} found but no languages defined. Using CLI args.")
    
    prs = Presentation(args.pptx)
    total_slides = len(prs.slides)
    
    # Initialize Gemini agent and runner
    logger.info("Initializing Gemini agent...")
    agent = create_agent()
    runner = InMemoryRunner(
        agent=agent,
        app_name='presentation_preloader',
    )
    
    db = firestore.Client(database="langbridge")

    print(f"\nProcessing {total_slides} slides from {args.pptx}")
    print(f"Languages: {', '.join(languages)}")
    if bucket_name:
        print(f"Speech bucket: {bucket_name}")
    else:
        print("Speech generation: DISABLED (no bucket configured)")
    print("Using Gemini to generate messages...\n")

    cached_count = 0
    speech_count = 0
    all_messages = {}

    # Process each slide individually
    for slide_idx in range(1, total_slides + 1):
        speaker_notes = get_speaker_notes(prs, slide_idx)
        
        if not speaker_notes:
            print(f"Slide {slide_idx}: No speaker notes - skipping")
            continue
        
        print(f"\nSlide {slide_idx}: Processing {len(speaker_notes)} "
              f"chars of notes")
        
        # Generate message for each language for this slide
        for lang in languages:
            logger.info(
                "Generating message for slide %d, language %s",
                slide_idx,
                lang
            )
            
            generated_message = generate_message_from_notes(
                runner,
                speaker_notes,
                lang,
                slide_idx
            )
            
            if not generated_message:
                logger.warning(
                    "Failed to generate message for slide %d, lang %s",
                    slide_idx,
                    lang
                )
                continue
            
            # Cache with speaker notes only (no slide number)
            # This way cache works even if slides are reordered
            cache_message(
                db,
                lang,
                speaker_notes,  # Just the notes content
                generated_message,
                course_id=args.course_id
            )
            cache_key = _cache_key(lang, speaker_notes)
            cached_count += 1
            
            # Store first 80 chars for display
            if len(generated_message) > 80:
                preview = generated_message[:80] + "..."
            else:
                preview = generated_message
            print(f"  [{lang}] Cached '{cache_key}'")
            print(f"       -> {preview}")
            
            # Generate speech file if bucket specified
            if bucket_name:
                try:
                    # Use course_utils to get voice params
                    voice_params = course_utils.get_voice_params(args.course_id, lang)
                    
                    speech_file = generate_speech_file(
                        bucket_name,
                        generated_message,
                        lang,
                        voice_params=voice_params
                    )
                    print(f"       -> Speech: {speech_file}")
                    speech_count += 1
                except Exception as e:
                    logger.warning(
                        "Failed to generate speech for slide %d, lang %s: %s",
                        slide_idx,
                        lang,
                        e
                    )
            
            # Build messages dict for config (keyed by lang:slideN)
            all_messages[f"{lang}:slide{slide_idx}"] = generated_message

    # Update config document
    if all_messages:
        update_config(db, all_messages)
        print(
            f"\n✓ Successfully cached {cached_count} AI-generated "
            f"message(s) across {len(languages)} language(s)"
        )
        if bucket_name and speech_count > 0:
            print(
                f"✓ Generated {speech_count} speech file(s) "
                f"in bucket '{bucket_name}'"
            )
        print(
            f"✓ Updated presentation_messages in config with "
            f"{len(all_messages)} entry/entries"
        )
    else:
        print("\n⚠ No messages generated")


if __name__ == "__main__":
    main()
